# build_pptx.py — Generator PPTX Sidang TA Pertanian Cerdas Cabai Rawit
# Jalankan: python build_pptx.py

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import copy

# ── Konstanta warna (dark theme) ────────────────────────────────
BG       = RGBColor(0x0f, 0x11, 0x17)
SLIDE_BG = RGBColor(0x1a, 0x1d, 0x27)
CARD     = RGBColor(0x22, 0x26, 0x37)
BORDER   = RGBColor(0x2e, 0x33, 0x50)
ACCENT   = RGBColor(0x4f, 0x8e, 0xf7)
GREEN    = RGBColor(0x22, 0xc5, 0x5e)
YELLOW   = RGBColor(0xf5, 0x9e, 0x0b)
RED      = RGBColor(0xef, 0x44, 0x44)
PURPLE   = RGBColor(0xa7, 0x8b, 0xfa)
TEAL     = RGBColor(0x34, 0xd3, 0x99)
WHITE    = RGBColor(0xe2, 0xe8, 0xf0)
MUTED    = RGBColor(0x94, 0xa3, 0xb8)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

# ── Helper: slide kosong dengan background gelap ─────────────────
def new_slide():
    blank = prs.slide_layouts[6]   # blank
    sl = prs.slides.add_slide(blank)
    bg = sl.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = SLIDE_BG
    return sl

# ── Helper: tambah textbox ────────────────────────────────────────
def tb(sl, text, x, y, w, h,
       size=18, bold=False, color=WHITE,
       align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = sl.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

# ── Helper: rectangle shape ───────────────────────────────────────
def rect(sl, x, y, w, h, fill_color=CARD, line_color=BORDER, line_w=Pt(1)):
    from pptx.util import Pt as oPt
    shape = sl.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color
    shape.line.width = line_w
    return shape

# ── Helper: top gradient bar (accent strip) ───────────────────────
def accent_bar(sl):
    bar = rect(sl, 0, 0, W, Pt(6), fill_color=ACCENT, line_color=ACCENT, line_w=Pt(0))
    return bar

# ── Helper: chip label ────────────────────────────────────────────
def chip(sl, text, x, y, color=ACCENT):
    r = rect(sl, x, y, Inches(3.2), Inches(0.32), fill_color=CARD, line_color=color, line_w=Pt(1))
    tb(sl, text, x + Inches(0.1), y + Inches(0.03), Inches(3.0), Inches(0.28),
       size=9, bold=True, color=color, align=PP_ALIGN.CENTER)

# ── Helper: section header ────────────────────────────────────────
def heading(sl, text, x=Inches(0.4), y=Inches(0.55)):
    tb(sl, text, x, y, Inches(12.5), Inches(0.7), size=26, bold=True, color=ACCENT)

# ── Helper: slide number (bottom right) ───────────────────────────
def slide_num(sl, n, total=12):
    tb(sl, f"{n} / {total}", Inches(12.1), Inches(7.1), Inches(1.0), Inches(0.3),
       size=9, color=BORDER, align=PP_ALIGN.RIGHT)

# ── Helper: bulleted list into a textbox ──────────────────────────
def bullets(sl, items, x, y, w, h, size=13, color=WHITE, bullet="  ✦  "):
    txBox = sl.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = bullet + item
        run.font.size = Pt(size)
        run.font.color.rgb = color

# ── Helper: card box with title and bullets ───────────────────────
def card_box(sl, title, items, x, y, w, h, title_color=PURPLE, size=12):
    rect(sl, x, y, w, h, fill_color=CARD, line_color=BORDER)
    tb(sl, title, x + Inches(0.15), y + Inches(0.1), w - Inches(0.3), Inches(0.35),
       size=13, bold=True, color=title_color)
    bullets(sl, items, x + Inches(0.1), y + Inches(0.45),
            w - Inches(0.2), h - Inches(0.55), size=size, color=WHITE)

# ── Helper: stat box (big number + label) ────────────────────────
def stat_box(sl, num, label, x, y, w=Inches(2.2), h=Inches(1.4),
             num_color=ACCENT):
    rect(sl, x, y, w, h)
    tb(sl, num, x, y + Inches(0.2), w, Inches(0.7),
       size=36, bold=True, color=num_color, align=PP_ALIGN.CENTER)
    tb(sl, label, x, y + Inches(0.9), w, Inches(0.4),
       size=11, color=MUTED, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
#  SLIDE 1 – COVER
# ════════════════════════════════════════════════════════════════
sl = new_slide()
accent_bar(sl)
slide_num(sl, 1)

chip(sl, "TUGAS AKHIR – TEKNIK INFORMATIKA", Inches(4.5), Inches(0.55))

tb(sl, "SISTEM PERTANIAN CERDAS\nTANAMAN CABE RAWIT",
   Inches(1.0), Inches(1.1), Inches(11.3), Inches(1.5),
   size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

tb(sl, "MENGGUNAKAN METODE FUZZY TAHANI BERBASIS IoT",
   Inches(1.0), Inches(2.6), Inches(11.3), Inches(0.7),
   size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# Info box
rect(sl, Inches(2.5), Inches(3.4), Inches(8.3), Inches(1.5))
tb(sl, "MUKHAMAD SA'ID AQIL  |  NIM. 2021150081",
   Inches(2.6), Inches(3.5), Inches(8.1), Inches(0.45),
   size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(sl, "Program Studi Teknik Informatika  ·  Fakultas Teknik dan Ilmu Komputer\n"
       "Universitas Sains Al-Qur'an Jawa Tengah di Wonosobo  ·  2026",
   Inches(2.6), Inches(3.95), Inches(8.1), Inches(0.8),
   size=11, color=MUTED, align=PP_ALIGN.CENTER)

# Pembimbing
rect(sl, Inches(0.5), Inches(5.1), Inches(6.0), Inches(0.55), line_color=ACCENT)
tb(sl, "Pembimbing Utama: Muslim Hidayat, S.Kom., M.Kom.",
   Inches(0.6), Inches(5.17), Inches(5.8), Inches(0.4),
   size=11, color=ACCENT, align=PP_ALIGN.CENTER)
rect(sl, Inches(6.8), Inches(5.1), Inches(6.0), Inches(0.55), line_color=ACCENT)
tb(sl, "Pembimbing Pendamping: Iman Ahmad Ihsannuddin, S.Pd.Kom, M.Pd.",
   Inches(6.9), Inches(5.17), Inches(5.8), Inches(0.4),
   size=11, color=ACCENT, align=PP_ALIGN.CENTER)

tb(sl, "🌶", Inches(6.2), Inches(0.45), Inches(1.0), Inches(0.8),
   size=36, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
#  SLIDE 2 – LATAR BELAKANG
# ════════════════════════════════════════════════════════════════
sl = new_slide()
accent_bar(sl)
slide_num(sl, 2)
chip(sl, "BAB I – PENDAHULUAN", Inches(0.4), Inches(0.1), ACCENT)
heading(sl, "Latar Belakang Masalah")

card_box(sl, "🌶  Urgensi Komoditas", [
    "Cabai rawit (Capsicum frutescens L.) bernilai ekonomi tinggi & strategis",
    "Kab. Magelang: lahan 3.828,9 ha — produksi 511.556,9 kwintal (BPS, 2023)",
    "Permintaan tinggi: rumah tangga & industri kuliner",
    "Kondisi ideal: kelembapan 50–70%, pH 6–7, suhu 24–28°C",
], Inches(0.4), Inches(1.4), Inches(6.0), Inches(2.8), title_color=GREEN)

card_box(sl, "⚠  Permasalahan Lapangan", [
    "Pemantauan masih MANUAL — petani harus datang ke greenhouse",
    "pH tanah, suhu, kelembapan sering tidak stabil",
    "Iklim lembap → risiko serangan jamur & OPT meningkat",
    "Kerontokan bunga → penurunan produktivitas",
    "Sistem monitoring tidak memiliki pengambilan keputusan otomatis",
], Inches(6.8), Inches(1.4), Inches(6.1), Inches(2.8), title_color=YELLOW)

rect(sl, Inches(0.4), Inches(4.35), Inches(12.5), Inches(0.85),
     fill_color=RGBColor(0x0f,0x1a,0x2e), line_color=ACCENT)
tb(sl, "💡  Solusi: Sistem Pertanian Cerdas berbasis IoT + Fuzzy Tahani untuk pemantauan "
       "real-time dan pengendalian otomatis kondisi lingkungan greenhouse cabai rawit.",
   Inches(0.55), Inches(4.42), Inches(12.2), Inches(0.7),
   size=13, bold=False, color=WHITE)

# ════════════════════════════════════════════════════════════════
#  SLIDE 3 – RUMUSAN, BATASAN, TUJUAN
# ════════════════════════════════════════════════════════════════
sl = new_slide()
accent_bar(sl)
slide_num(sl, 3)
chip(sl, "BAB I – PENDAHULUAN", Inches(0.4), Inches(0.1), PURPLE)
heading(sl, "Rumusan Masalah, Batasan & Tujuan")

card_box(sl, "❓  Rumusan Masalah", [
    "Bagaimana merancang sistem IoT yang mampu memantau dan mengendalikan kondisi cabai rawit secara real-time?",
    "Bagaimana menerapkan Fuzzy Tahani + MQTT untuk menghasilkan keputusan pengendalian otomatis?",
], Inches(0.4), Inches(1.4), Inches(4.0), Inches(2.7), title_color=ACCENT)

card_box(sl, "🔒  Batasan Masalah", [
    "Objek: cabai rawit, Kab. Magelang, skala greenhouse",
    "3 sensor: DHT22, Soil Moisture, pH Tanah",
    "Data dikirim via protokol MQTT",
    "Ditampilkan pada dashboard web real-time",
], Inches(4.65), Inches(1.4), Inches(4.0), Inches(2.7), title_color=YELLOW)

card_box(sl, "🎯  Tujuan Penelitian", [
    "Membangun sistem IoT (ESP32) untuk monitoring pH, suhu, kelembapan secara real-time",
    "Menerapkan Fuzzy Tahani + MQTT untuk keputusan otomatis yang cepat & stabil",
], Inches(8.9), Inches(1.4), Inches(4.0), Inches(2.7), title_color=GREEN)

# Manfaat
rect(sl, Inches(0.4), Inches(4.3), Inches(12.5), Inches(2.85))
tb(sl, "Manfaat Penelitian", Inches(0.55), Inches(4.38), Inches(5.0), Inches(0.35),
   size=13, bold=True, color=PURPLE)
bullets(sl, [
    "Pengguna: monitoring jarak jauh via MQTT, hemat air & energi, kurangi risiko kerusakan tanaman",
    "IPTEK: referensi sistem IoT + MQTT untuk pertanian cerdas, contoh implementasi Fuzzy Tahani",
    "Peneliti: pemahaman integrasi MQTT, IoT, Fuzzy Tahani dalam satu sistem",
], Inches(0.55), Inches(4.75), Inches(12.0), Inches(2.2), size=12)

# ════════════════════════════════════════════════════════════════
#  SLIDE 4 – LANDASAN TEORI
# ════════════════════════════════════════════════════════════════
sl = new_slide()
accent_bar(sl)
slide_num(sl, 4)
chip(sl, "BAB II – LANDASAN TEORI", Inches(0.4), Inches(0.1), GREEN)
heading(sl, "Landasan Teori")

theories = [
    ("🌐  IoT", "Jaringan objek fisik terhubung internet untuk akuisisi data, monitoring & pengendalian otomatis.", ACCENT),
    ("🔧  ESP32", "Mikrokontroler dual-core 240 MHz, WiFi+BT terintegrasi, 520 KB SRAM. Pusat pengolahan data IoT.", PURPLE),
    ("📡  MQTT", "Protokol ringan publish-subscribe. Latensi rendah, efisien bandwidth, cocok perangkat berdaya rendah.", GREEN),
    ("🧠  Fuzzy Tahani", "Metode keputusan fuzzy berbasis DB relasional. Query fuzzy + fire strength → keputusan adaptif.", YELLOW),
]
for i, (title, desc, col) in enumerate(theories):
    x = Inches(0.4 + i * 3.22)
    rect(sl, x, Inches(1.35), Inches(3.1), Inches(2.0), line_color=col)
    tb(sl, title, x + Inches(0.12), Inches(1.42), Inches(2.9), Inches(0.45),
       size=13, bold=True, color=col)
    tb(sl, desc, x + Inches(0.12), Inches(1.85), Inches(2.9), Inches(1.35),
       size=11, color=WHITE)

sensors = [
    ("🌡  DHT22", "Suhu udara −40~80°C & kelembapan 0–100% RH", ACCENT),
    ("🪱  Soil Moisture", "Kadar air tanah via resistansi/kapasitansi. Dasar keputusan penyiraman otomatis.", GREEN),
    ("⚗  Sensor pH Tanah", "Keasaman tanah. Cabai rawit: pH optimal 6–7 untuk penyerapan hara efektif.", YELLOW),
    ("🔌  Relay & Pompa", "Relay: saklar aktuator. Pompa Air: irigasi. Pompa pH: koreksi keasaman. Kipas: suhu.", PURPLE),
]
for i, (title, desc, col) in enumerate(sensors):
    x = Inches(0.4 + i * 3.22)
    rect(sl, x, Inches(3.55), Inches(3.1), Inches(2.0), line_color=col)
    tb(sl, title, x + Inches(0.12), Inches(3.62), Inches(2.9), Inches(0.45),
       size=13, bold=True, color=col)
    tb(sl, desc, x + Inches(0.12), Inches(4.05), Inches(2.9), Inches(1.35),
       size=11, color=WHITE)

# ════════════════════════════════════════════════════════════════
#  SLIDE 5 – METODOLOGI
# ════════════════════════════════════════════════════════════════
sl = new_slide()
accent_bar(sl)
slide_num(sl, 5)
chip(sl, "BAB III – METODOLOGI PENELITIAN", Inches(0.4), Inches(0.1), YELLOW)
heading(sl, "Metodologi Penelitian (Prototyping 5 Tahap)")

# Kiri: flow prototyping
steps = [
    ("1", "Komunikasi", "Wawancara 3 petani + observasi lapangan", ACCENT),
    ("2", "Pengumpulan Kebutuhan", "Analisis hardware & software sistem", PURPLE),
    ("3", "Membangun Sistem", "Diagram blok + skema rangkaian", TEAL),
    ("4", "Pengkodean", "C++ ESP32 + HTML/CSS/JS Dashboard", YELLOW),
    ("5", "Pengujian", "Black Box Testing + User Acceptance Test", RED),
]
for i, (num, title, desc, col) in enumerate(steps):
    y = Inches(1.35 + i * 1.12)
    rect(sl, Inches(0.4), y, Inches(6.0), Inches(1.0), line_color=col)
    tb(sl, num, Inches(0.5), y + Inches(0.1), Inches(0.4), Inches(0.8),
       size=22, bold=True, color=col, align=PP_ALIGN.CENTER)
    tb(sl, title, Inches(1.0), y + Inches(0.05), Inches(5.2), Inches(0.38),
       size=13, bold=True, color=col)
    tb(sl, desc, Inches(1.0), y + Inches(0.45), Inches(5.2), Inches(0.45),
       size=11, color=MUTED)
    if i < 4:
        tb(sl, "▼", Inches(3.1), y + Inches(1.02), Inches(0.5), Inches(0.2),
           size=10, color=ACCENT, align=PP_ALIGN.CENTER)

# Kanan: tabel himpunan fuzzy
rect(sl, Inches(6.8), Inches(1.35), Inches(6.1), Inches(5.75))
tb(sl, "Variabel Input & Himpunan Fuzzy", Inches(6.95), Inches(1.42),
   Inches(5.9), Inches(0.4), size=13, bold=True, color=PURPLE)

rows = [
    ("Suhu Udara", "Rendah", "≤ 24°C", ACCENT),
    ("", "Sedang", "24 – 31°C", ACCENT),
    ("", "Tinggi", "≥ 27°C", ACCENT),
    ("Kelembapan Tanah", "Kering", "≤ 40%", GREEN),
    ("", "Lembab", "40 – 80%", GREEN),
    ("", "Basah", "≥ 70%", GREEN),
    ("pH Tanah", "Asam", "3 – 6", YELLOW),
    ("", "Normal", "5.5 – 7.5", YELLOW),
    ("", "Basa", "7 – 9", YELLOW),
]
headers = ["Variabel", "Himpunan", "Range"]
for j, h in enumerate(headers):
    x = Inches(6.95 + j * 1.95)
    rect(sl, x, Inches(1.87), Inches(1.9), Inches(0.35),
         fill_color=RGBColor(0x1a,0x1d,0x27), line_color=BORDER)
    tb(sl, h, x + Inches(0.05), Inches(1.9), Inches(1.8), Inches(0.3),
       size=11, bold=True, color=MUTED)
for i, (var, hmp, rng, col) in enumerate(rows):
    y = Inches(2.25 + i * 0.5)
    for j, (val, c) in enumerate([(var, col), (hmp, WHITE), (rng, MUTED)]):
        x = Inches(6.95 + j * 1.95)
        tb(sl, val, x + Inches(0.05), y + Inches(0.05),
           Inches(1.8), Inches(0.4), size=11, color=c)

# ════════════════════════════════════════════════════════════════
#  SLIDE 6 – DESAIN SISTEM
# ════════════════════════════════════════════════════════════════
sl = new_slide()
accent_bar(sl)
slide_num(sl, 6)
chip(sl, "BAB IV – HASIL & PEMBAHASAN", Inches(0.4), Inches(0.1), ACCENT)
heading(sl, "Desain Sistem – Arsitektur & Hardware")

# Kiri: diagram blok
tb(sl, "Diagram Blok Sistem", Inches(0.4), Inches(1.35), Inches(5.8), Inches(0.38),
   size=13, bold=True, color=PURPLE)

sensor_boxes = [("🌡 DHT22","Suhu+RH"), ("🪱 Soil Moisture","Kelembapan"), ("⚗ pH Sensor","pH Tanah")]
for i, (lbl, sub) in enumerate(sensor_boxes):
    x = Inches(0.4 + i * 1.92)
    rect(sl, x, Inches(1.8), Inches(1.8), Inches(0.75), line_color=ACCENT)
    tb(sl, lbl,  x+Inches(0.05), Inches(1.85), Inches(1.7), Inches(0.38), size=11, bold=True, color=ACCENT)
    tb(sl, sub, x+Inches(0.05), Inches(2.22), Inches(1.7), Inches(0.28), size=10, color=MUTED)

tb(sl, "▼▼▼", Inches(2.5), Inches(2.6), Inches(1.5), Inches(0.3), size=12, color=ACCENT, align=PP_ALIGN.CENTER)

rect(sl, Inches(0.4), Inches(2.95), Inches(5.8), Inches(0.95), line_color=PURPLE)
tb(sl, "⚙  ESP32 — FreeRTOS Dual-Core", Inches(0.55), Inches(3.0), Inches(5.5), Inches(0.38),
   size=13, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
tb(sl, "Core 0: MQTT + Supabase  |  Core 1: Sensor + Fuzzy + Relay",
   Inches(0.55), Inches(3.38), Inches(5.5), Inches(0.45), size=10, color=MUTED, align=PP_ALIGN.CENTER)

aktuator = [("🌀 Kipas","Suhu Tinggi"), ("💦 Pompa Air","Tanah Kering"), ("🧪 Pompa pH","pH Asam")]
for i, (lbl, sub) in enumerate(aktuator):
    x = Inches(0.4 + i * 1.92)
    rect(sl, x, Inches(4.05), Inches(1.8), Inches(0.75), line_color=GREEN)
    tb(sl, lbl,  x+Inches(0.05), Inches(4.1),  Inches(1.7), Inches(0.38), size=11, bold=True, color=GREEN)
    tb(sl, sub, x+Inches(0.05), Inches(4.47), Inches(1.7), Inches(0.28), size=10, color=MUTED)

tb(sl, "▼", Inches(2.8), Inches(4.85), Inches(0.7), Inches(0.3), size=12, color=ACCENT, align=PP_ALIGN.CENTER)

rect(sl, Inches(0.4), Inches(5.2), Inches(2.75), Inches(0.65), line_color=YELLOW)
tb(sl, "📡 MQTT (EMQX)", Inches(0.5), Inches(5.27), Inches(2.55), Inches(0.22), size=11, bold=True, color=YELLOW)
tb(sl, "TLS 8883 / WSS 8084", Inches(0.5), Inches(5.49), Inches(2.55), Inches(0.25), size=9, color=MUTED)
rect(sl, Inches(3.45), Inches(5.2), Inches(2.75), Inches(0.65), line_color=YELLOW)
tb(sl, "🗄 Supabase DB", Inches(3.55), Inches(5.27), Inches(2.55), Inches(0.22), size=11, bold=True, color=YELLOW)
tb(sl, "REST API / HTTPS", Inches(3.55), Inches(5.49), Inches(2.55), Inches(0.25), size=9, color=MUTED)

tb(sl, "▼", Inches(2.8), Inches(5.9), Inches(0.7), Inches(0.3), size=12, color=ACCENT, align=PP_ALIGN.CENTER)
rect(sl, Inches(0.4), Inches(6.25), Inches(5.8), Inches(0.65), line_color=RED)
tb(sl, "🖥  Dashboard Web Real-time + Kontrol Manual",
   Inches(0.55), Inches(6.32), Inches(5.5), Inches(0.45), size=12, bold=True, color=RED, align=PP_ALIGN.CENTER)

# Kanan: tabel HW & SW
tb(sl, "Kebutuhan Hardware & Software", Inches(7.0), Inches(1.35), Inches(5.9), Inches(0.38),
   size=13, bold=True, color=PURPLE)
hw_sw = [
    ("[HW]", "ESP32",         "Mikrokontroler utama", ACCENT),
    ("[HW]", "DHT22",         "Suhu & kelembapan udara", ACCENT),
    ("[HW]", "Soil Moisture", "Kelembapan tanah", ACCENT),
    ("[HW]", "Sensor pH+DMS", "pH tanah + modul driver", ACCENT),
    ("[HW]", "Relay 1+2ch",  "Saklar aktuator", ACCENT),
    ("[HW]", "2 Pompa Mini",  "Irigasi & koreksi pH", ACCENT),
    ("[SW]", "Arduino IDE",   "Pemrograman ESP32", GREEN),
    ("[SW]", "EMQX Cloud",    "MQTT Broker TLS", GREEN),
    ("[SW]", "Supabase",      "Database cloud REST", GREEN),
    ("[SW]", "VS Code + HTML","Dashboard web", GREEN),
]
for i, (cat, name, func, col) in enumerate(hw_sw):
    y = Inches(1.8 + i * 0.5)
    rect(sl, Inches(7.0), y, Inches(5.9), Inches(0.46), line_color=BORDER)
    tb(sl, cat,  Inches(7.05), y+Inches(0.05), Inches(0.6),  Inches(0.36), size=10, bold=True, color=col)
    tb(sl, name, Inches(7.7),  y+Inches(0.05), Inches(1.8),  Inches(0.36), size=11, bold=True, color=WHITE)
    tb(sl, func, Inches(9.55), y+Inches(0.05), Inches(3.2),  Inches(0.36), size=11, color=MUTED)

# ════════════════════════════════════════════════════════════════
#  SLIDE 7 – IMPLEMENTASI FUZZY TAHANI
# ════════════════════════════════════════════════════════════════
sl = new_slide()
accent_bar(sl)
slide_num(sl, 7)
chip(sl, "BAB IV – HASIL & PEMBAHASAN", Inches(0.4), Inches(0.1), PURPLE)
heading(sl, "Implementasi Fuzzy Tahani & Arsitektur FreeRTOS")

# Kiri: rule base
tb(sl, "Rule Base Fuzzy Tahani", Inches(0.4), Inches(1.35), Inches(5.8), Inches(0.38),
   size=13, bold=True, color=PURPLE)
rules = [
    ("R1", "IF Suhu TINGGI",    "μtinggi(T)", "🌀 Kipas ON",    ACCENT),
    ("R2", "IF Tanah KERING",   "μkering(S)", "💦 Pompa Air ON", GREEN),
    ("R3", "IF pH ASAM",        "μasam(P)",   "🧪 Pompa pH ON",  YELLOW),
]
hdrs = ["Rule", "Kondisi", "Fire Strength", "Aktuator"]
for j, h in enumerate(hdrs):
    xw = [0.5, 1.5, 1.8, 2.0]
    x = Inches(0.4 + sum(xw[:j]))
    rect(sl, x, Inches(1.8), Inches(xw[j]), Inches(0.4),
         fill_color=CARD, line_color=BORDER)
    tb(sl, h, x+Inches(0.05), Inches(1.83), Inches(xw[j]-0.1), Inches(0.32),
       size=11, bold=True, color=MUTED)
for i, (r, cond, fs, act, col) in enumerate(rules):
    y = Inches(2.23 + i * 0.6)
    vals = [r, cond, fs, act]
    xw = [0.5, 1.5, 1.8, 2.0]
    for j, (v, c) in enumerate(zip(vals, [col, WHITE, MUTED, GREEN])):
        x = Inches(0.4 + sum(xw[:j]))
        rect(sl, x, y, Inches(xw[j]), Inches(0.55), line_color=BORDER)
        tb(sl, v, x+Inches(0.05), y+Inches(0.07), Inches(xw[j]-0.1), Inches(0.4),
           size=12, bold=(j==0), color=col if j==0 else c)

rect(sl, Inches(0.4), Inches(4.1), Inches(5.8), Inches(1.3), line_color=PURPLE)
tb(sl, "Threshold Keputusan:", Inches(0.55), Inches(4.17), Inches(5.5), Inches(0.35),
   size=12, bold=True, color=PURPLE)
bullets(sl, [
    "Kipas ON jika μtinggi > 0.5",
    "Pompa Air ON jika μkering > 0.4  dan  soil > 0",
    "Pompa pH ON jika μasam > 0.4  dan  pH > 0",
], Inches(0.55), Inches(4.5), Inches(5.5), Inches(0.85), size=11)

rect(sl, Inches(0.4), Inches(5.5), Inches(5.8), Inches(1.35), line_color=GREEN)
tb(sl, "Logika Timer Pompa (Anti-boros):", Inches(0.55), Inches(5.57), Inches(5.5), Inches(0.35),
   size=12, bold=True, color=GREEN)
bullets(sl, [
    "Pompa Air: nyala 10 detik → jeda 30 menit (auto)",
    "Pompa pH: nyala 10 detik → jeda 3 jam (auto)",
    "Manual mode: bypass semua jeda — langsung nyala",
], Inches(0.55), Inches(5.97), Inches(5.5), Inches(0.75), size=11)

# Kanan: FreeRTOS
tb(sl, "Arsitektur FreeRTOS Dual-Core ESP32", Inches(6.8), Inches(1.35), Inches(6.1), Inches(0.38),
   size=13, bold=True, color=TEAL)

rect(sl, Inches(6.8), Inches(1.8), Inches(6.1), Inches(2.15), line_color=ACCENT)
tb(sl, "Core 0 — TaskKomunikasi (Prioritas 1)", Inches(6.95), Inches(1.87),
   Inches(5.8), Inches(0.38), size=12, bold=True, color=ACCENT)
bullets(sl, [
    "WiFi connect + NTP time sync",
    "MQTT TLS keep-alive + callback perintah manual",
    "Publish data sensor ke broker tiap 15 detik",
    "Insert Supabase REST API tiap 60 detik",
], Inches(6.95), Inches(2.28), Inches(5.8), Inches(1.55), size=11)

rect(sl, Inches(6.8), Inches(4.1), Inches(6.1), Inches(2.15), line_color=TEAL)
tb(sl, "Core 1 — TaskSensor (Prioritas 2)", Inches(6.95), Inches(4.17),
   Inches(5.8), Inches(0.38), size=12, bold=True, color=TEAL)
bullets(sl, [
    "Baca DHT22, Soil (Kalman Filter), pH (Kalman + drift-guard)",
    "Inferensi Fuzzy Tahani → keputusan aktuator",
    "tickPompa() tiap 20ms — timer pompa auto & manual",
    "Fast-path relay manual < 20ms tanpa tunggu siklus sensor",
], Inches(6.95), Inches(4.58), Inches(5.8), Inches(1.55), size=11)

rect(sl, Inches(6.8), Inches(6.35), Inches(6.1), Inches(0.6), line_color=YELLOW)
tb(sl, "Sinkronisasi: xSemaphoreMutex lindungi SharedData antar core  ·  volatile ManualCmd untuk perintah instan",
   Inches(6.95), Inches(6.42), Inches(5.8), Inches(0.45), size=10, color=YELLOW)

# ════════════════════════════════════════════════════════════════
#  SLIDE 8 – TAMPILAN SISTEM
# ════════════════════════════════════════════════════════════════
sl = new_slide()
accent_bar(sl)
slide_num(sl, 8)
chip(sl, "BAB IV – HASIL & PEMBAHASAN", Inches(0.4), Inches(0.1), ACCENT)
heading(sl, "Tampilan Sistem – Dashboard Web & Prototype")

# Kiri: mockup
tb(sl, "Dashboard Web Real-Time", Inches(0.4), Inches(1.35), Inches(6.5), Inches(0.38),
   size=13, bold=True, color=PURPLE)

# header bar mockup
rect(sl, Inches(0.4), Inches(1.8), Inches(6.5), Inches(0.5),
     fill_color=RGBColor(0x1a,0x1d,0x27), line_color=BORDER)
tb(sl, "🌶 Pertanian Cerdas Cabai Rawit", Inches(0.5), Inches(1.85),
   Inches(4.5), Inches(0.38), size=11, bold=True, color=ACCENT)
rect(sl, Inches(5.7), Inches(1.88), Inches(1.1), Inches(0.32),
     fill_color=RGBColor(0x16,0x65,0x34), line_color=GREEN)
tb(sl, "● LIVE", Inches(5.72), Inches(1.9), Inches(1.0), Inches(0.28),
   size=9, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

cards_mock = [("SUHU","24.0°C","Rendah",YELLOW), ("TANAH","81.3%","Lembab",GREEN),
              ("pH","6.90","Normal",GREEN), ("HUMID","80%","–",MUTED)]
for i, (lbl, val, badge, col) in enumerate(cards_mock):
    x = Inches(0.4 + i * 1.6)
    rect(sl, x, Inches(2.38), Inches(1.52), Inches(1.1), line_color=col)
    tb(sl, lbl, x+Inches(0.05), Inches(2.42), Inches(1.42), Inches(0.28), size=9, color=MUTED)
    tb(sl, val, x+Inches(0.05), Inches(2.68), Inches(1.42), Inches(0.45), size=14, bold=True, color=col)
    tb(sl, badge, x+Inches(0.05), Inches(3.1), Inches(1.42), Inches(0.3), size=9, color=WHITE)

relay_mock = [("🌀 Kipas","OFF",MUTED), ("💦 Pompa Air","ON",GREEN), ("🧪 Pompa pH","OFF",MUTED)]
for i, (lbl, st, col) in enumerate(relay_mock):
    x = Inches(0.4 + i * 2.15)
    rect(sl, x, Inches(3.6), Inches(2.0), Inches(0.75), line_color=col)
    tb(sl, lbl, x+Inches(0.05), Inches(3.65), Inches(1.9), Inches(0.3), size=10, bold=True, color=WHITE)
    tb(sl, st, x+Inches(0.05), Inches(3.98), Inches(1.9), Inches(0.3), size=12, bold=True, color=col)

rect(sl, Inches(0.4), Inches(4.45), Inches(6.5), Inches(0.55), line_color=BORDER)
tb(sl, "📋 Riwayat 120 baris terakhir dari Supabase — auto-refresh 60 detik",
   Inches(0.5), Inches(4.52), Inches(6.3), Inches(0.38), size=10, color=MUTED)

# Kanan: fitur
tb(sl, "Fitur Utama Dashboard", Inches(7.2), Inches(1.35), Inches(5.7), Inches(0.38),
   size=13, bold=True, color=PURPLE)
bullets(sl, [
    "Tampilan real-time: suhu, RH, kelembapan tanah, pH — update tiap 15 detik via MQTT",
    "Badge status himpunan fuzzy: Rendah/Sedang/Tinggi · Kering/Lembab/Basah · Asam/Normal/Basa",
    "Indikator relay ON/OFF + countdown jeda pompa (detik tersisa)",
    "Kontrol manual aktuator via dashboard → kirim perintah via MQTT topic pertanian/kontrol",
    "Riwayat data Supabase (120 baris) — auto-refresh & reload manual",
    "Status WiFi ESP32 (SSID, IP, RSSI) tampil di header",
], Inches(7.2), Inches(1.8), Inches(5.7), Inches(3.2), size=12)

rect(sl, Inches(7.2), Inches(5.15), Inches(5.7), Inches(1.75))
tb(sl, "Stack Teknologi", Inches(7.35), Inches(5.22), Inches(5.4), Inches(0.38),
   size=12, bold=True, color=TEAL)
bullets(sl, [
    "HTML + CSS dark theme · MQTT.js WebSocket",
    "Supabase REST API untuk penyimpanan historis",
    "Deployed via GitHub Pages (akses publik)",
], Inches(7.35), Inches(5.62), Inches(5.4), Inches(1.2), size=11)
